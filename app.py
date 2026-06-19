from flask import Flask, render_template, request, redirect, session, jsonify, send_file, send_from_directory, Response, stream_with_context
import os
import base64
import re
import json
import tempfile
import sqlite3
import time
import requests
from io import BytesIO
from dotenv import load_dotenv
from authlib.integrations.flask_client import OAuth
from groq import Groq
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
import pdfplumber
from PIL import Image
from deepsearch import run_deepsearch

load_dotenv()

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 16 * 1024 * 1024
app.config.update(
    SESSION_COOKIE_SAMESITE="Lax",
    SESSION_COOKIE_SECURE=os.getenv("FLASK_ENV") == "production",
    SESSION_COOKIE_HTTPONLY=True,
)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "dev-secret-key-change-me")

oauth = OAuth(app)
oauth.register(
    name="google",
    client_id=os.getenv("GOOGLE_CLIENT_ID"),
    client_secret=os.getenv("GOOGLE_CLIENT_SECRET"),
    server_metadata_url="https://accounts.google.com/.well-known/openid-configuration",
    client_kwargs={"scope": "openid email profile", "prompt": "select_account"},
)

client = Groq(api_key=os.getenv("GROQ_API_KEY"))
CLOUDFLARE_API_TOKEN  = os.getenv("CLOUDFLARE_API_TOKEN", "")
CLOUDFLARE_ACCOUNT_ID = os.getenv("CLOUDFLARE_ACCOUNT_ID", "")

# ═══════════════════════════════════════════════════════
#  DATABASE
# ═══════════════════════════════════════════════════════

def get_db():
    conn = sqlite3.connect("pranox.db", check_same_thread=False)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    db = get_db()
    db.execute("""CREATE TABLE IF NOT EXISTS chats(
        id INTEGER PRIMARY KEY AUTOINCREMENT, user_email TEXT, role TEXT,
        message TEXT, created_at DATETIME DEFAULT CURRENT_TIMESTAMP)""")
    db.execute("""CREATE TABLE IF NOT EXISTS user_memory(
        id INTEGER PRIMARY KEY AUTOINCREMENT, user_email TEXT, key TEXT, value TEXT)""")
    db.commit()
    db.close()

init_db()

def save_memory(user_email, key, value):
    db = get_db()
    db.execute("DELETE FROM user_memory WHERE user_email=? AND key=?", (user_email, key))
    db.execute("INSERT INTO user_memory(user_email,key,value) VALUES (?,?,?)", (user_email, key, value))
    db.commit()
    db.close()

def get_memory(user_email):
    db = get_db()
    rows = db.execute("SELECT key,value FROM user_memory WHERE user_email=?", (user_email,)).fetchall()
    db.close()
    return "\n".join([f"{r['key']}: {r['value']}" for r in rows])


# ═══════════════════════════════════════════════════════
#  AI-DRIVEN SEARCH DECISION
#  Uses llama-3.1-8b-instant (fast, cheap) to decide:
#  1. Does this question need a live web search?
#  2. What is the best Google query to use?
#  This handles ALL topics automatically — no keyword lists needed.
# ═══════════════════════════════════════════════════════

def ai_search_decision(user_message: str) -> tuple:
    """
    Returns: (needs_search: bool, search_query: str)
    """
    prompt = f"""You are a search decision engine for an AI assistant. Your only job is to decide:
1. Does the user's question need a live web search to answer correctly?
2. If yes, write the single best Google search query to find the answer.

SEARCH IS REQUIRED for:
- Current or recent news, events, updates, announcements
- People's current roles or positions (Chief Minister, PM, President, CEO, Minister, Owner, Director, etc.)
- Prices of anything (gold, fuel, crypto, stocks, products, services)
- Sports scores, match results, standings, schedules
- Weather forecasts or current conditions
- Recently launched or released products, apps, movies, shows, songs
- Election results, political changes, government news
- Company news, acquisitions, funding, leadership changes
- Location of any business, shop, restaurant, hospital, branch, office, mall
- Whether a specific business exists in a city ("is there a Zara in Mysore")
- "Who won", "what happened", "latest", "recent", "current", "now", "today"
- Anything that changes over time and the model might have outdated info

SEARCH IS NOT REQUIRED for:
- Math, calculations, unit conversions
- Coding, programming, technical concepts
- Creative writing: stories, poems, essays, emails
- Definitions, meanings, translations
- Historical facts that cannot change
- Grammar, spelling, language questions
- How things work (science, tech explanations that don't change)
- Personal advice, opinions, recommendations based on user preferences
- Questions about Pranox AI itself

User question: "{user_message}"

Reply in EXACTLY this format. No extra text, no explanation:
SEARCH: YES or NO
QUERY: the best Google search query if YES, or NONE if NO"""

    try:
        completion = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.0,
            max_tokens=80,
        )
        response = completion.choices[0].message.content.strip()
        lines = response.splitlines()

        search_line = next((l for l in lines if l.upper().startswith("SEARCH:")), "SEARCH: NO")
        query_line  = next((l for l in lines if l.upper().startswith("QUERY:")),  "QUERY: NONE")

        needs_search = "YES" in search_line.upper()
        query = query_line.split(":", 1)[1].strip() if ":" in query_line else ""
        query = "" if query.upper() == "NONE" else query

        if needs_search and not query:
            query = user_message  # fallback

        print(f"[AI DECISION] Search={needs_search} | Query='{query}'")
        return needs_search, query

    except Exception as e:
        print(f"[AI DECISION ERROR] {e} — defaulting to search")
        return True, user_message  # on error, always search to be safe


# ═══════════════════════════════════════════════════════
#  SERPER WEB SEARCH
# ═══════════════════════════════════════════════════════

def search_internet(query: str) -> str:
    try:
        api_key = os.getenv("SERPER_API_KEY")
        if not api_key:
            print("[SEARCH] SERPER_API_KEY not set")
            return ""

        res = requests.post(
            "https://google.serper.dev/search",
            headers={"X-API-KEY": api_key, "Content-Type": "application/json"},
            json={"q": query, "num": 8, "gl": "in", "hl": "en"},
            timeout=8,
        )

        if res.status_code != 200:
            print(f"[SEARCH] Serper returned {res.status_code}")
            return ""

        data = res.json()
        results = []

        # 1. Answer box — highest priority, direct answer
        if "answerBox" in data:
            ab = data["answerBox"]
            answer = ab.get("answer") or ab.get("snippet") or ab.get("title", "")
            source = ab.get("link", "")
            if answer:
                results.append(
                    f"[TOP ANSWER]: {answer}" + (f"\n  Source: {source}" if source else "")
                )

        # 2. Knowledge graph — structured facts
        if "knowledgeGraph" in data:
            kg    = data["knowledgeGraph"]
            title = kg.get("title", "")
            desc  = kg.get("description", "")
            attrs = kg.get("attributes", {})
            if title or desc:
                kg_text = f"[KNOWLEDGE GRAPH]: {title}"
                if desc:
                    kg_text += f" — {desc}"
                for k, v in list(attrs.items())[:6]:
                    kg_text += f"\n  {k}: {v}"
                results.append(kg_text)

        # 3. Organic results — web pages and news
        if "organic" in data:
            for r in data["organic"][:6]:
                snippet = r.get("snippet", "")
                title   = r.get("title", "")
                link    = r.get("link", "")
                if snippet:
                    results.append(f"[SOURCE: {title}]\n  {snippet}\n  Link: {link}")

        # 4. Related searches — only if very few results
        if "relatedSearches" in data and len(results) < 3:
            related = [r.get("query", "") for r in data["relatedSearches"][:3]]
            if related:
                results.append(f"[Related searches]: {', '.join(related)}")

        final = "\n\n".join(results)
        print(f"[SEARCH SUCCESS] {len(results)} results for: '{query}'")
        return final

    except requests.Timeout:
        print("[SEARCH] Timeout after 8s")
        return ""
    except Exception as e:
        print(f"[SEARCH ERROR]: {e}")
        return ""


# ═══════════════════════════════════════════════════════
#  UTILITIES
# ═══════════════════════════════════════════════════════

def safe_trim(text, limit=6000):
    return text[:limit] if len(text) > limit else text

def clean_text(text):
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()

MODELS = ["llama-3.3-70b-versatile", "llama-3.1-8b-instant"]

def run_ai(messages, model_index=0, max_tokens=1200):
    # ── ROLLBACK / RETRY LOGIC ────────────────────────────
    # On transient crashes (rate limits, timeouts, 5xx, overload)
    # the same model is retried up to MAX_RETRIES times with
    # exponential backoff before rolling back to the next model.
    # Permanent errors (bad request, auth) skip retries immediately.
    MAX_RETRIES  = 2          # attempts per model on transient errors
    BASE_DELAY   = 1.5        # seconds (multiplied by attempt number)
    TRANSIENT    = ("rate", "timeout", "429", "500", "502", "503",
                    "overload", "connection", "unavailable", "reset")

    for i in range(model_index, len(MODELS)):
        for attempt in range(MAX_RETRIES + 1):
            try:
                completion = client.chat.completions.create(
                    model=MODELS[i],
                    messages=messages,
                    temperature=0.3,
                    max_tokens=max_tokens,
                )
                return completion.choices[0].message.content.strip()
            except Exception as e:
                err_str = str(e).lower()
                is_transient = any(kw in err_str for kw in TRANSIENT)
                if is_transient and attempt < MAX_RETRIES:
                    wait = BASE_DELAY * (attempt + 1)
                    print(f"[ROLLBACK] Transient crash on {MODELS[i]} "
                          f"(attempt {attempt+1}/{MAX_RETRIES}): {e} — retrying in {wait}s")
                    time.sleep(wait)
                    continue          # retry same model
                # Permanent error or retries exhausted → roll to next model
                print(f"AI ERROR (model {MODELS[i]}) after {attempt+1} attempt(s):", e)
                break
    return None
    # ─────────────────────────────────────────────────────

# ═══════════════════════════════════════════════════════
#  THINKING EFFORT — user-selectable response depth
#  Purely controls model choice + answer length, scaling
#  monotonically from low -> medium -> high. The VISIBLE
#  "thinking" (live reasoning shown to the user) is a separate
#  concern handled by /api/chat/stream + REASONING_TOKEN_BUDGET
#  below, so the two controls no longer fight each other.
# ═══════════════════════════════════════════════════════

EFFORT_CONFIG = {
    "off":    {"model_index": 0, "max_tokens": 900},
    "low":    {"model_index": 1, "max_tokens": 500},
    "medium": {"model_index": 0, "max_tokens": 1300},
    "high":   {"model_index": 0, "max_tokens": 2200},
}

# How many tokens the live "thinking" pass gets, scaled by effort.
REASONING_TOKEN_BUDGET = {
    "off": 250, "low": 200, "medium": 350, "high": 550,
}

# How many discrete reasoning steps to ask for, scaled by effort.
REASONING_STEP_COUNT = {
    "off": "2-4", "low": "2-3", "medium": "4-6", "high": "6-9",
}

# Explicit depth/length instruction injected into the final-answer prompt so
# "high" effort reliably reads as more thorough than "low" — not just a
# bigger token cap, but an actual instruction to go deeper.
EFFORT_DEPTH_INSTRUCTIONS = {
    "off":    "Give a clear, direct answer.",
    "low":    "Keep the answer brief and to the point — just the essentials, no filler.",
    "medium": "Give a clear, well-rounded answer with the key points explained.",
    "high":   "Give a deep, thorough answer: cover the topic fully, explain the reasoning "
              "behind it, and include relevant detail, nuance, or examples where useful.",
}

def normalize_effort(value):
    value = (value or "off").strip().lower()
    return value if value in EFFORT_CONFIG else "off"

def run_ai_with_effort(messages, effort="off"):
    cfg = EFFORT_CONFIG.get(effort, EFFORT_CONFIG["off"])
    return run_ai(messages, model_index=cfg["model_index"], max_tokens=cfg["max_tokens"])

def is_bad_response(reply):
    if not reply:
        return True
    bad_patterns = [
        "here's the corrected code",
        "flask application",
        "example of how you could",
        "missing code",
    ]
    return any(p in reply.lower() for p in bad_patterns)

def is_hedging_response(reply):
    """Detect if the model ignored search results and fell back to stale training-data hedging."""
    if not reply:
        return True
    hedging_phrases = [
        "my data may be outdated",
        "as of my last update",
        "my knowledge cutoff",
        "i don't have real-time",
        "i cannot confirm",
        "my training data",
        "as of my training",
        "i suggest verifying",
        "may not be accurate",
        "i cannot provide real-time",
        "as of my knowledge",
        "i don't have access to real",
    ]
    reply_lower = reply.lower()
    return any(p in reply_lower for p in hedging_phrases)


# ═══════════════════════════════════════════════════════
#  SYSTEM PROMPTS
# ═══════════════════════════════════════════════════════

BASE_SYSTEM_PROMPT = """You are Pranox AI — a next-generation AI assistant. You are highly intelligent, friendly, and precise — built to work like the best AI assistants in the world.

════════════════════════════
IDENTITY
════════════════════════════
- Your name: Pranox AI
- Founder: Chetansaipranav R
- Created: January 2026
- If anyone asks "who created you", "who is your founder", "tell me about Pranox":
  Always answer: "Chetansaipranav R is the founder of Pranox AI. He created it in January 2026."

════════════════════════════
THINKING (INTERNAL — NEVER SHOW TO USER)
════════════════════════════
Before every response:
- Understand what the user is truly asking
- Think through the answer step by step internally
- Verify the answer makes sense before writing it
- For live search results: check that the result actually answers the question
- NEVER show your thinking or reasoning steps to the user
- Give only the final clean, correct answer

════════════════════════════
GREETING RULES
════════════════════════════
When user says: hi, hello, hey, hii, hiya, good morning, good afternoon, good evening, sup, what's up:
- Respond warmly and briefly
- If user's name is in memory, include it naturally
- Examples:
  "Hi! How can I help you today?"
  "Hello Pranav! What can I do for you?"
  "Hey there! Need any help?"
- NEVER say "bye" in response to a greeting
- NEVER give an unrelated or weird response to a greeting
- If user repeats a greeting, respond politely every single time
- Keep it short, warm, and natural — never robotic

════════════════════════════
PERSONALIZATION
════════════════════════════
- If user's name is in memory, use it naturally — mainly in greetings or the opening line
- Do NOT repeat the name in every sentence — use it once at most
- When user tells you their name for the first time: "Nice to meet you, [name]!"
- Keep personalization natural and human-feeling

<<<<<<< HEAD
════════════════════════════
CONVERSATION & CONTEXT
════════════════════════════
- Always use the conversation history to maintain continuity
- Never repeat the same answer unnecessarily
- If user refers to something from earlier in the chat, connect it correctly
=======
========================
THINKING (IMPORTANT)
========================
- Understand the question deeply
- Break into logical steps internally
- Do NOT show reasoning
- Give only final clean answer
- Before answering think and check that while it is correct or not internally 
- Before giving response about latest information or frequently changing information through browsing check if it is correct or not internally
- While user is asking about the questions like general knowledge, like constantly changing things for example weather, about politics, company CEO's , owners of companies and etc, always browse and check it twice internally is it correct or not and give the response correctly
>>>>>>> 49db481 (updated)

════════════════════════════
CODE RULES
════════════════════════════
- Always use fenced code blocks with correct language syntax highlighting
- Clean indentation — no sloppy formatting
- Separate explanation from code clearly
- Briefly explain what the code does before or after the block

════════════════════════════
RESPONSE FORMAT
<<<<<<< HEAD
════════════════════════════
1. Give the direct answer or main point first
2. Use bullet points only when listing multiple distinct items
3. Clean spacing between sections — never walls of text
4. Code goes in proper fenced code blocks
5. End with 1–2 follow-up question suggestions when they genuinely help the user go deeper
   (Skip follow-ups for: simple greetings, yes/no answers, very casual one-liners)
=======
========================
1. Clear explanation first
2. Bullets when needed
3. Clean spacing
4. Code in proper code blocks
5. Always ask follow-up questions 
6. When users ask about latest news or anything about recent happened browse and give the answer
7. When providing a response think and check while it is correct or not internally while giving the response of latest information from browsing
8. If you are not sure about the answer or you think it is a bad response, do not give that response and instead search on the internet and give the correct answer based on that search results
9. If you are giving the response based on browsing or search results, always check twice internally while giving the response that it is correct or not based on those search results and then give the response
10. Always give response correctly by browsing the internet when user is asking about the questions like general knowledge, like constantly changing things for example weather, about politics, company CEO's , owners of companies and etc
>>>>>>> 49db481 (updated)

════════════════════════════
BEHAVIOR
════════════════════════════
- Be friendly, smart, and helpful — like ChatGPT or Claude
- Never be robotic or give stiff corporate-sounding responses
- Keep answers structured and easy to read
- For short casual questions: give short casual answers
- For complex questions: give thorough, well-organized answers

════════════════════════════
PRANOX LINKS
════════════════════════════
Share ONLY when user explicitly asks about Pranox AI's social media or official pages:
- Instagram: https://www.instagram.com/pranoxgroups
- X (Twitter): https://x.com/Pranoxgroups
- LinkedIn: https://www.linkedin.com/in/chetansaipranav-r-a6b18333b
- Product Hunt: https://www.producthunt.com/@pranoxai
NEVER share these links in unrelated conversations.

════════════════════════════
CONTACT / EMAIL
════════════════════════════
Share ONLY when user explicitly asks for contact or email info:
- Email: pranoxoffical@gmail.com
"""

SEARCH_OVERRIDE_PROMPT = """
╔══════════════════════════════════════════════════╗
║   🔴 LIVE SEARCH MODE — READ THIS CAREFULLY 🔴  ║
╚══════════════════════════════════════════════════╝

Real-time web search results have been fetched RIGHT NOW for this question.
These results are from TODAY. They are the ground truth. They override your training data.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
STEP 1 — READ THE RESULTS
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
- Check [TOP ANSWER] first — if the answer is there, use it directly.
- Check [KNOWLEDGE GRAPH] next for structured facts.
- Check [SOURCE] snippets for additional context and details.
- The answer IS in the results. Read carefully before responding.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
STEP 2 — ANSWER DIRECTLY AND CONFIDENTLY
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
State the answer in the very first sentence. Examples:
✅ "The Chief Minister of West Bengal is Mamata Banerjee."
✅ "Gold price today is ₹7,200 per gram."
✅ "The latest iPhone model is the iPhone 16 Pro."
✅ "McDonald's in Davangere is located at P.J. Extension."

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
STEP 3 — BANNED PHRASES (NEVER WRITE ANY OF THESE)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
These phrases are COMPLETELY FORBIDDEN. Writing them is a failure:
❌ "my data may be outdated"
❌ "as of my last update"
❌ "as of my knowledge cutoff"
❌ "as of my training"
❌ "I don't have real-time data"
❌ "I don't have real-time access"
❌ "I cannot confirm"
❌ "I cannot verify"
❌ "please verify this information"
❌ "I suggest verifying"
❌ "may not be accurate"
❌ "my training data"
❌ "served as" (when the person CURRENTLY holds the position)

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
STEP 4 — IF RESULTS CONFLICT
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Prefer: Wikipedia > official government/company sites > major news outlets > YouTube
Say: "Most sources indicate X." — this is allowed.
Never pick a random YouTube video over Wikipedia as the truth.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
STEP 5 — IF RESULTS ARE GENUINELY EMPTY
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Only if the search results section below is actually empty or has no useful data, say:
"I searched but couldn't find specific details. Try Google Maps or the official website for the most accurate info."
╔══════════════════════════════════════════════════╗
║              SEARCH RESULTS BELOW               ║
╚══════════════════════════════════════════════════╝
"""

def build_system_prompt(memory: str) -> str:
    prompt  = BASE_SYSTEM_PROMPT
    prompt += "\n════════════════════════════\nUSER MEMORY\n════════════════════════════\n"
    prompt += f"{memory}\n" if memory else "No memory stored yet.\n"
    return prompt

def build_messages_with_search(system_prompt, search_results, history, user_message):
    msgs = [{"role": "system", "content": system_prompt}]
    # Add conversation history oldest-first, skipping history[0] (current user message)
    for h in reversed(list(history)[1:]):
        role = h["role"] if h["role"] in ("user", "assistant") else "user"
        msgs.append({"role": role, "content": h["message"]})
    # Inject search results directly into the user turn — the model cannot ignore this
    enriched = (
        f"{SEARCH_OVERRIDE_PROMPT}\n"
        f"{search_results}\n"
        f"{'═'*50}\n\n"
        f"USER QUESTION: {user_message}\n\n"
        f"Answer the USER QUESTION using the search results above. "
        f"Be direct and confident. Give the answer in your first sentence."
    )
    msgs.append({"role": "user", "content": enriched})
    return msgs

def build_messages_no_search(system_prompt, history, user_message):
    msgs = [{"role": "system", "content": system_prompt}]
    for h in reversed(list(history)[1:]):
        role = h["role"] if h["role"] in ("user", "assistant") else "user"
        msgs.append({"role": role, "content": h["message"]})
    msgs.append({"role": "user", "content": user_message})
    return msgs


def get_redirect_uri():
    render_url = os.getenv("RENDER_EXTERNAL_URL")
    if render_url:
        return f"{render_url.rstrip('/')}/authorize"
    return "http://127.0.0.1:8000/authorize"


# ═══════════════════════════════════════════════════════
#  ROUTES — AUTH
# ═══════════════════════════════════════════════════════

@app.route("/")
def landing():
    return render_template("landing.html")

@app.route("/login")
def login():
    redirect_uri = get_redirect_uri()
    print("REDIRECT URI:", redirect_uri)
    return oauth.google.authorize_redirect(redirect_uri)

@app.route("/authorize")
def authorize():
    try:
        token = oauth.google.authorize_access_token()
        user  = oauth.google.userinfo()
        session["user"] = {
            "email":   user.get("email"),
            "name":    user.get("name"),
            "picture": user.get("picture"),
        }
        return redirect("/dashboard")
    except Exception as e:
        print("GOOGLE LOGIN ERROR:", e)
        return redirect("/?error=login_failed")

@app.route("/logout")
def logout():
    session.clear()
    return redirect("/")

@app.route("/dashboard")
def dashboard():
    if "user" not in session:
        return redirect("/login")
    return render_template("dashboard.html", user=session.get("user"))

@app.route("/chat")
def chat():
    return render_template("chat.html", user=session.get("user"))


# ═══════════════════════════════════════════════════════
#  ROUTES — TOOLS
# ═══════════════════════════════════════════════════════

@app.route("/email", methods=["GET", "POST"])
def email():
    if "user" not in session:
        return render_template("login_required.html")
    output = ""
    if request.method == "POST":
        topic  = request.form.get("topic", "").strip()
        tone   = request.form.get("tone", "Professional").strip()
        length = request.form.get("length", "Medium").strip()

        length_guide = {
            "Short":  "Write a concise email of 3-4 short paragraphs (around 100-150 words).",
            "Medium": "Write a well-developed email of 4-6 paragraphs (around 200-300 words).",
            "Long":   "Write a detailed, thorough email of 6-8 paragraphs (around 350-500 words) covering every aspect the user mentioned.",
        }.get(length, "Write a well-developed email of 4-6 paragraphs.")

        system_prompt = (
            "You are an expert professional email writer with 15+ years of experience writing emails "
            "for executives, businesses, and individuals across all industries. "
            "Your emails are always clear, polished, and highly effective. "
            "RULES:\n"
            "- Write a COMPLETE, READY-TO-SEND email — include Subject line, greeting, body, and sign-off.\n"
            "- Use the exact tone and intent described by the user. Understand their requirement deeply and write accordingly.\n"
            "- Structure: Subject: [subject] then blank line then greeting then body paragraphs then closing then signature placeholder.\n"
            "- NO markdown symbols (* # _ ` ~). Plain text only.\n"
            "- Do NOT add meta-commentary like Here is your email or I hope this helps. Just write the email.\n"
            "- Make the email sound natural, human, and purposeful — never generic or template-like.\n"
            "- Every paragraph must add real value and move the email forward."
        )

        user_prompt = (
            f"Write a {tone.lower()} email for the following requirement:\n\n"
            f"{topic}\n\n"
            f"{length_guide}\n"
            f"Ensure the email fully addresses everything the user described. "
            f"Make it compelling, accurate, and complete."
        )

        output = run_ai([
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_prompt},
        ], max_tokens=1800) or "Couldn't generate email. Please try again."
        output = clean_text(re.sub(r"[*#_`~]", "", output))
    return render_template("email.html", email=output)

# ═══════════════════════════════════════════════════════
#  RESUME ROUTE  —  REPLACE THIS BLOCK IN app.py
#  Only this route is changed. Nothing else is touched.
# ═══════════════════════════════════════════════════════

@app.route("/resume", methods=["GET", "POST"])
def resume():
    if "user" not in session:
        return render_template("login_required.html")
    output = ""
    if request.method == "POST":
        name       = request.form.get("name", "").strip()
        role       = request.form.get("role", "").strip()
        skills     = request.form.get("skills", "").strip()
        experience = request.form.get("experience", "").strip()
        education  = request.form.get("education", "").strip()

        if not all([name, role, skills, experience, education]):
            return render_template("resume.html", resume="Please fill all fields!")

        prompt = (
            f"Create a professional resume for:\n"
            f"Name: {name}\n"
            f"Target Role: {role}\n"
            f"Skills: {skills}\n"
            f"Work Experience: {experience}\n"
            f"Education: {education}\n\n"
            f"Format the resume using EXACTLY this structure:\n"
            f"{name.upper()}\n"
            f"{role}\n\n"
            f"PROFESSIONAL SUMMARY\n"
            f"[Write 2-3 sentence professional summary]\n\n"
            f"SKILLS\n"
            f"[List each skill on its own line starting with a bullet •]\n\n"
            f"WORK EXPERIENCE\n"
            f"[Job Title] | [Company] | [Duration]\n"
            f"[List achievements starting with bullet •]\n\n"
            f"EDUCATION\n"
            f"[Degree] | [Institution] | [Year]\n"
            f"[One sentence about education if needed]\n"
        )

        output = run_ai(
            [
                {
                    "role": "system",
                    "content": (
                        "You are an expert professional resume writer. "
                        "Write ATS-optimized, interview-winning resumes. "
                        "Follow the exact format given by the user precisely. "
                        "Use ALL CAPS only for section headers: PROFESSIONAL SUMMARY, SKILLS, WORK EXPERIENCE, EDUCATION. "
                        "Use the bullet character • (not asterisks or dashes) for all list items. "
                        "Use pipe | to separate job title, company, and dates. "
                        "Write clean, confident, action-verb-driven bullet points. "
                        "Do NOT use asterisks *, hashes #, underscores _, or backticks `. "
                        "Do NOT use markdown. Plain structured text only."
                    ),
                },
                {"role": "user", "content": prompt},
            ],
            max_tokens=1500,
        ) or "Couldn't generate resume. Please try again."

        # Strip any stray markdown symbols but preserve • bullets and | separators
        output = re.sub(r"[*#_`]", "", output)
        output = re.sub(r"\n{3,}", "\n\n", output).strip()

    return render_template("resume.html", resume=output)

# ═══════════════════════════════════════════════════════
#  MAIN CHAT API
# ═══════════════════════════════════════════════════════

@app.route("/api/chat", methods=["POST"])
def api_chat():
    data         = request.get_json(force=True)
    user_message = safe_trim(data.get("message", "").strip())
    effort       = normalize_effort(data.get("effort"))
    if not user_message:
        return jsonify({"reply": "Please send a message."}), 400

    user_email = session["user"]["email"] if "user" in session else "guest"
    db = get_db()
    try:
        # Save user message to DB
        db.execute(
            "INSERT INTO chats(user_email,role,message) VALUES (?,?,?)",
            (user_email, "user", user_message)
        )
        db.commit()

        # Load memory
        memory = get_memory(user_email)

        # Auto-save name from message
        name_match = re.search(r"my name is ([a-zA-Z ]+)", user_message, re.IGNORECASE)
        if name_match:
            save_memory(user_email, "name", name_match.group(1).strip().title())
            memory = get_memory(user_email)

        # Auto-save age from message
        age_match = re.search(r"i(?:'m| am) (\d+) years? old", user_message, re.IGNORECASE)
        if age_match:
            save_memory(user_email, "age", age_match.group(1))
            memory = get_memory(user_email)

        # Fetch recent conversation history (DESC = newest first, index 0 = current message)
        history = db.execute(
            "SELECT role,message FROM chats WHERE user_email=? ORDER BY id DESC LIMIT 20",
            (user_email,)
        ).fetchall()
        history = [h for h in history if "couldn't fully process" not in h["message"].lower()]

        # ── AI-DRIVEN SEARCH DECISION ─────────────────────────
        # The AI model decides if search is needed and writes the best query.
        # This handles ANY topic — no hardcoded keyword lists required.
        do_search, search_query = ai_search_decision(user_message)
        search_results = ""

        if do_search and search_query:
            search_results = search_internet(search_query)
            # Fallback: if the AI-optimized query returned nothing, try the raw message
            if not search_results:
                print(f"[FALLBACK SEARCH] Trying raw: '{user_message}'")
                search_results = search_internet(user_message)
        # ─────────────────────────────────────────────────────

        system_prompt = build_system_prompt(memory)

        # Build messages based on whether search results are available
        if search_results.strip():
            msgs = build_messages_with_search(system_prompt, search_results, history, user_message)
        elif do_search and not search_results:
            # Search was triggered but Serper returned nothing
            no_result_note = (
                f"NOTE: A live web search was attempted for this question but returned no results. "
                f"Answer from your training knowledge as best you can. "
                f"For location-specific or very recent questions, suggest the user "
                f"check Google Maps or the official website.\n\n"
                f"USER QUESTION: {user_message}"
            )
            msgs = [{"role": "system", "content": system_prompt}]
            for h in reversed(list(history)[1:]):
                role = h["role"] if h["role"] in ("user", "assistant") else "user"
                msgs.append({"role": role, "content": h["message"]})
            msgs.append({"role": "user", "content": no_result_note})
        else:
            msgs = build_messages_no_search(system_prompt, history, user_message)

        # Get AI reply
        reply = run_ai_with_effort(msgs, effort)
        if not reply:
            reply = "I ran into an issue generating a response. Please try again."

        # ── HEDGING DETECTION + RETRY ─────────────────────────
        # If the model ignored search results and hedged with stale phrases, retry once
        # with a minimal stripped-down prompt that forces it to use the search data.
        if search_results.strip() and is_hedging_response(reply):
            print("[RETRY] Model hedged despite search results — retrying with direct prompt")
            retry_msgs = [
                {
                    "role": "system",
                    "content": (
                        "You are a factual assistant. You have live search results. "
                        "Answer ONLY from those results. "
                        "NEVER write 'my data may be outdated', 'as of my last update', "
                        "'I don't have real-time data', or any similar hedging phrase. "
                        "State the answer directly and confidently in 1-2 sentences."
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        f"Search results:\n{search_results}\n\n"
                        f"Question: {user_message}\n\n"
                        f"Give a direct answer from the search results above. "
                        f"Start with the answer immediately. No hedging."
                    ),
                },
            ]
            retry_reply = run_ai(retry_msgs)
            if retry_reply and not is_hedging_response(retry_reply):
                reply = retry_reply
                print("[RETRY] Retry succeeded")
            else:
                print("[RETRY] Retry also hedged — keeping original reply")
        # ─────────────────────────────────────────────────────

        # Catch completely broken/irrelevant responses
        if is_bad_response(reply):
            fallback = search_internet(user_message)
            reply = (
                f"Here are some helpful resources:\n\n{fallback}"
                if fallback
                else "I couldn't process that. Could you rephrase your question?"
            )

        reply = re.sub(r"\n{3,}", "\n\n", reply).strip()

        # Nudge guests to log in (once per session)
        if "user" not in session and "login_nudge" not in session:
            reply += "\n\n💡 *Login to save your chat history and get personalized responses*"
            session["login_nudge"] = True

        # Save assistant reply to DB
        if "couldn't fully process" not in reply.lower():
            db.execute(
                "INSERT INTO chats(user_email,role,message) VALUES (?,?,?)",
                (user_email, "assistant", reply)
            )
            db.commit()

        return jsonify({"reply": reply})

    except Exception as e:
        print("CHAT ERROR:", e)
        return jsonify({"reply": "An unexpected error occurred. Please try again."}), 500
    finally:
        db.close()


# ═══════════════════════════════════════════════════════
#  LIVE THINKING — real-time streamed reasoning + answer
#  Used when the "Thinking" toggle is ON in the UI.
#  Unlike /api/chat, this streams the model's ACTUAL
#  reasoning tokens (not a canned animation) over SSE,
#  followed by the actual final-answer tokens, both live.
# ═══════════════════════════════════════════════════════

def _build_chat_context(user_message):
    """Shared prep: memory, history, search decision, message list. Mirrors /api/chat."""
    user_email = session["user"]["email"] if "user" in session else "guest"
    db = get_db()
    try:
        db.execute(
            "INSERT INTO chats(user_email,role,message) VALUES (?,?,?)",
            (user_email, "user", user_message)
        )
        db.commit()

        memory = get_memory(user_email)

        name_match = re.search(r"my name is ([a-zA-Z ]+)", user_message, re.IGNORECASE)
        if name_match:
            save_memory(user_email, "name", name_match.group(1).strip().title())
            memory = get_memory(user_email)

        age_match = re.search(r"i(?:'m| am) (\d+) years? old", user_message, re.IGNORECASE)
        if age_match:
            save_memory(user_email, "age", age_match.group(1))
            memory = get_memory(user_email)

        history = db.execute(
            "SELECT role,message FROM chats WHERE user_email=? ORDER BY id DESC LIMIT 20",
            (user_email,)
        ).fetchall()
        history = [h for h in history if "couldn't fully process" not in h["message"].lower()]
    finally:
        db.close()

    do_search, search_query = ai_search_decision(user_message)
    search_results = ""
    if do_search and search_query:
        search_results = search_internet(search_query)
        if not search_results:
            search_results = search_internet(user_message)

    system_prompt = build_system_prompt(memory)

    if search_results.strip():
        msgs = build_messages_with_search(system_prompt, search_results, history, user_message)
    elif do_search and not search_results:
        no_result_note = (
            f"NOTE: A live web search was attempted for this question but returned no results. "
            f"Answer from your training knowledge as best you can. "
            f"For location-specific or very recent questions, suggest the user "
            f"check Google Maps or the official website.\n\n"
            f"USER QUESTION: {user_message}"
        )
        msgs = [{"role": "system", "content": system_prompt}]
        for h in reversed(list(history)[1:]):
            role = h["role"] if h["role"] in ("user", "assistant") else "user"
            msgs.append({"role": role, "content": h["message"]})
        msgs.append({"role": "user", "content": no_result_note})
    else:
        msgs = build_messages_no_search(system_prompt, history, user_message)

    return msgs, user_email


@app.route("/api/chat/stream", methods=["POST"])
def api_chat_stream():
    data         = request.get_json(force=True)
    user_message = safe_trim(data.get("message", "").strip())
    effort       = normalize_effort(data.get("effort"))
    if not user_message:
        return jsonify({"reply": "Please send a message."}), 400

    msgs, user_email = _build_chat_context(user_message)
    cfg = EFFORT_CONFIG.get(effort, EFFORT_CONFIG["off"])
    reasoning_tokens = REASONING_TOKEN_BUDGET.get(effort, 250)
    step_count = REASONING_STEP_COUNT.get(effort, "3-5")
    depth_note = EFFORT_DEPTH_INSTRUCTIONS.get(effort, EFFORT_DEPTH_INSTRUCTIONS["off"])

    # Decide the login nudge BEFORE streaming starts, since the session
    # cookie can only be updated while response headers are still being built.
    add_nudge = "user" not in session and "login_nudge" not in session
    if add_nudge:
        session["login_nudge"] = True

    def generate():
        last = msgs[-1]

        # ── PASS 1: genuine live reasoning, streamed token-by-token ──
        # Asked for one short step per line so the UI can render a real,
        # Claude-style step timeline instead of one big paragraph.
        reasoning_messages = msgs[:-1] + [{
            "role": "user",
            "content": (
                f"{last['content']}\n\n"
                f"Think this through step by step, like quick working notes. Write "
                f"{step_count} short steps, ONE PER LINE, no numbering or bullets — just "
                f"the plain text of each step (e.g. what's being asked, key facts or "
                f"context that matter, including any search results above, and how the "
                f"pieces fit together). Each line under 14 words. Do NOT write the final "
                f"answer here — only the reasoning steps, one per line."
            ),
        }]

        # Same transient-error retry policy as run_ai(), so a momentary
        # rate-limit/timeout/5xx from Groq doesn't kill the whole stream.
        STREAM_MAX_RETRIES = 2
        STREAM_BASE_DELAY  = 1.5
        STREAM_TRANSIENT   = ("rate", "timeout", "429", "500", "502", "503",
                               "overload", "connection", "unavailable", "reset")

        full_thinking = ""
        for attempt in range(STREAM_MAX_RETRIES + 1):
            try:
                stream = client.chat.completions.create(
                    model=MODELS[0],
                    messages=reasoning_messages,
                    temperature=0.4,
                    max_tokens=reasoning_tokens,
                    stream=True,
                )
                for chunk in stream:
                    delta = chunk.choices[0].delta.content or ""
                    if delta:
                        full_thinking += delta
                        yield f"event: thinking\ndata: {json.dumps(delta)}\n\n"
                break
            except Exception as e:
                print(f"THINKING STREAM ERROR (attempt {attempt+1}):", e)
                if full_thinking:
                    break  # partial content already streamed — don't retry mid-stream
                is_transient = any(kw in str(e).lower() for kw in STREAM_TRANSIENT)
                if is_transient and attempt < STREAM_MAX_RETRIES:
                    time.sleep(STREAM_BASE_DELAY * (attempt + 1))
                    continue
                break  # permanent error or retries exhausted — thinking is optional, move on

        yield f"event: thinking_done\ndata: {json.dumps({})}\n\n"

        # ── PASS 2: real final answer, informed by the reasoning above, streamed live ──
        # depth_note scales actual answer depth/length with effort, on top of the
        # max_tokens budget, so "high" reliably gives a fuller answer than "low".
        final_messages = msgs[:-1] + [{
            "role": "user",
            "content": (
                f"{last['content']}\n\n"
                + (f"Your reasoning above:\n{full_thinking}\n\n" if full_thinking else "")
                + f"{depth_note} "
                  "Do not repeat the reasoning steps — just the polished final answer."
            ),
        }]

        full_reply = ""
        for model_i in range(cfg["model_index"], len(MODELS)):
            if full_reply:
                break  # already got a reply from a previous model in this loop
            for attempt in range(STREAM_MAX_RETRIES + 1):
                try:
                    stream = client.chat.completions.create(
                        model=MODELS[model_i],
                        messages=final_messages,
                        temperature=0.3,
                        max_tokens=cfg["max_tokens"],
                        stream=True,
                    )
                    for chunk in stream:
                        delta = chunk.choices[0].delta.content or ""
                        if delta:
                            full_reply += delta
                            yield f"event: answer\ndata: {json.dumps(delta)}\n\n"
                    break
                except Exception as e:
                    print(f"ANSWER STREAM ERROR (model {MODELS[model_i]}, attempt {attempt+1}):", e)
                    if full_reply:
                        break  # partial content already streamed — don't retry mid-stream
                    is_transient = any(kw in str(e).lower() for kw in STREAM_TRANSIENT)
                    if is_transient and attempt < STREAM_MAX_RETRIES:
                        time.sleep(STREAM_BASE_DELAY * (attempt + 1))
                        continue
                    break  # permanent error or retries exhausted on this model -> roll to next model

        if not full_reply.strip():
            full_reply = "I ran into an issue generating a response. Please try again."

        full_reply = re.sub(r"\n{3,}", "\n\n", full_reply).strip()
        if add_nudge:
            full_reply += "\n\n💡 *Login to save your chat history and get personalized responses*"

        db2 = get_db()
        try:
            if "couldn't fully process" not in full_reply.lower():
                db2.execute(
                    "INSERT INTO chats(user_email,role,message) VALUES (?,?,?)",
                    (user_email, "assistant", full_reply)
                )
                db2.commit()
        finally:
            db2.close()

        yield f"event: done\ndata: {json.dumps({'reply': full_reply})}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


# ═══════════════════════════════════════════════════════
#  IMAGE VISION — Groq multimodal
# ═══════════════════════════════════════════════════════

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff", ".tif"}
IMAGE_MIME = {
    ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png",
    ".gif": "image/gif", ".webp": "image/webp",
    ".bmp": "image/png", ".tiff": "image/png", ".tif": "image/png",
}
VISION_MODELS = [
    "meta-llama/llama-4-scout-17b-16e-instruct",
    "meta-llama/llama-4-maverick-17b-128e-instruct",
]

def analyse_image(image_bytes, ext, user_question=""):
    mime = IMAGE_MIME.get(ext, "image/jpeg")
    if ext in (".bmp", ".tiff", ".tif"):
        try:
            img = Image.open(BytesIO(image_bytes)).convert("RGB")
            buf = BytesIO()
            img.save(buf, format="PNG")
            image_bytes = buf.getvalue()
            mime = "image/png"
        except Exception as e:
            print("Image conversion error:", e)
    b64      = base64.standard_b64encode(image_bytes).decode("utf-8")
    data_url = f"data:{mime};base64,{b64}"
    prompt   = (
        user_question.strip() if user_question.strip()
        else (
            "Describe this image in full detail. List every object, person, text, "
            "chart or scene visible. If there is readable text, transcribe it exactly."
        )
    )
    for model in VISION_MODELS:
        try:
            completion = client.chat.completions.create(
                model=model,
                messages=[{"role": "user", "content": [
                    {"type": "image_url", "image_url": {"url": data_url}},
                    {"type": "text", "text": prompt},
                ]}],
                max_tokens=1024,
                temperature=0.5,
            )
            return completion.choices[0].message.content.strip()
        except Exception as e:
            print(f"Vision error ({model}):", e)
    return ""


# ═══════════════════════════════════════════════════════
#  TEXT EXTRACTION from uploaded files
# ═══════════════════════════════════════════════════════

def extract_text(file_bytes, ext):
    text = ""

    if ext == ".pdf":
        try:
            with pdfplumber.open(BytesIO(file_bytes)) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"
        except Exception as e:
            print("PDF error:", e)

    elif ext in (".docx", ".doc"):
        try:
            from docx import Document
            doc   = Document(BytesIO(file_bytes))
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    rt = "\t".join(c.text for c in row.cells if c.text.strip())
                    if rt.strip():
                        parts.append(rt)
            text = "\n".join(parts)
        except Exception as e:
            print("python-docx error:", e)
            try:
                import docx2txt
                with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                text = docx2txt.process(tmp_path)
                os.unlink(tmp_path)
            except Exception as e2:
                print("docx2txt error:", e2)

    elif ext in (".xlsx", ".xlsm"):
        try:
            import openpyxl
            wb    = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
            parts = []
            for sn in wb.sheetnames:
                ws = wb[sn]
                parts.append(f"[Sheet: {sn}]")
                for row in ws.iter_rows(values_only=True):
                    rs = "\t".join(str(c) if c is not None else "" for c in row)
                    if rs.strip():
                        parts.append(rs)
            text = "\n".join(parts)
        except Exception as e:
            print("openpyxl error:", e)

    elif ext == ".xls":
        try:
            import xlrd
            wb    = xlrd.open_workbook(file_contents=file_bytes)
            parts = []
            for sheet in wb.sheets():
                parts.append(f"[Sheet: {sheet.name}]")
                for ri in range(sheet.nrows):
                    rs = "\t".join(str(sheet.cell_value(ri, ci)) for ci in range(sheet.ncols))
                    if rs.strip():
                        parts.append(rs)
            text = "\n".join(parts)
        except Exception as e:
            print("xlrd error:", e)

    elif ext == ".ods":
        try:
            import pandas as pd
            with tempfile.NamedTemporaryFile(delete=False, suffix=".ods") as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            df_dict = pd.read_excel(tmp_path, engine="odf", sheet_name=None)
            os.unlink(tmp_path)
            parts = []
            for sheet, df in df_dict.items():
                parts.append(f"[Sheet: {sheet}]")
                parts.append(df.to_string(index=False))
            text = "\n\n".join(parts)
        except Exception as e:
            print("ODS error:", e)

    elif ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs   = Presentation(BytesIO(file_bytes))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            text = "\n".join(parts)
        except Exception as e:
            print("PPTX error:", e)

    elif ext == ".rtf":
        try:
            from striprtf.striprtf import rtf_to_text
            text = rtf_to_text(file_bytes.decode("utf-8", errors="ignore"))
        except Exception as e:
            print("RTF error:", e)

    elif ext in (".csv", ".tsv"):
        try:
            import pandas as pd
            sep = "\t" if ext == ".tsv" else ","
            df  = pd.read_csv(BytesIO(file_bytes), sep=sep)
            text = df.to_string(index=False)
        except Exception as e:
            print("CSV error:", e)
            text = file_bytes.decode("utf-8", errors="ignore")

    elif ext in (".json", ".jsonl"):
        try:
            raw  = file_bytes.decode("utf-8", errors="ignore")
            text = (
                json.dumps(json.loads(raw), indent=2)
                if ext == ".json"
                else "\n".join(raw.splitlines()[:50])
            )
        except Exception as e:
            print("JSON error:", e)
            text = file_bytes.decode("utf-8", errors="ignore")[:4000]

    elif ext == ".ipynb":
        try:
            nb    = json.loads(file_bytes.decode("utf-8", errors="ignore"))
            parts = []
            for cell in nb.get("cells", []):
                src = "".join(cell.get("source", []))
                if src.strip():
                    parts.append(f"[{cell.get('cell_type','').upper()}]\n{src}")
            text = "\n\n".join(parts)
        except Exception as e:
            print("IPYNB error:", e)

    elif ext == ".odt":
        try:
            import odf.opendocument
            with tempfile.NamedTemporaryFile(delete=False, suffix=".odt") as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            doc = odf.opendocument.load(tmp_path)
            os.unlink(tmp_path)
            parts = []
            for el in doc.text.childNodes:
                s = el.plaintext() if hasattr(el, "plaintext") else str(el)
                if s.strip():
                    parts.append(s.strip())
            text = "\n".join(parts)
        except Exception as e:
            print("ODT error:", e)

    elif ext == ".xml":
        text = "XML Data:\n" + file_bytes.decode("utf-8", errors="ignore")

    elif ext == ".zip":
        try:
            import zipfile
            parts = []
            with zipfile.ZipFile(BytesIO(file_bytes)) as zf:
                for name in zf.namelist():
                    if name.endswith((
                        ".txt", ".md", ".csv", ".json", ".xml",
                        ".py", ".js", ".html", ".css", ".yaml", ".yml", ".sql"
                    )):
                        try:
                            content = zf.read(name).decode("utf-8", errors="ignore")
                            parts.append(f"[File: {name}]\n{content}")
                        except Exception:
                            pass
            text = "\n\n".join(parts)
        except Exception as e:
            print("ZIP error:", e)

    else:
        text = file_bytes.decode("utf-8", errors="ignore")

    return text.strip()


# ═══════════════════════════════════════════════════════
#  FILE UPLOAD ROUTE — multi-file (up to MAX_UPLOAD_FILES)
#  Reads every uploaded file (images + documents) and answers
#  using ALL of them together in a single synthesised reply,
#  the same way ChatGPT/Claude handle multi-file uploads.
# ═══════════════════════════════════════════════════════

MAX_UPLOAD_FILES     = 10     # max files accepted per message
MAX_CHARS_PER_FILE   = 3000   # extracted-text budget per file
MAX_TOTAL_FILE_CHARS = 20000  # combined extracted-text budget for the AI call

@app.route("/api/upload", methods=["POST"])
def upload():
    # Accept multiple files sent under "file" (input multiple / repeated append)
    # or "files", whichever the frontend uses — without changing anything else.
    files = request.files.getlist("file") or request.files.getlist("files")
    files = [f for f in files if f and f.filename]

    if not files:
        return jsonify({"reply": "No file received. Please try again."})

    if len(files) > MAX_UPLOAD_FILES:
        files = files[:MAX_UPLOAD_FILES]

    user_question = request.form.get("message", "").strip()
    effort        = normalize_effort(request.form.get("effort"))
    user_email    = session["user"]["email"] if "user" in session else "guest"

    file_sections   = []   # extracted content fed to the model, one block per file
    processed_names = []   # filenames successfully read
    failed_names    = []   # filenames that could not be read/extracted

    for file in files:
        filename = file.filename
        _, ext   = os.path.splitext(filename.lower())

        try:
            file_bytes = file.read()
        except Exception as e:
            print(f"FILE READ ERROR [{filename}]:", e)
            failed_names.append(filename)
            continue

        if ext in IMAGE_EXTENSIONS:
            description = analyse_image(file_bytes, ext, user_question)
            if description:
                file_sections.append(f"[Image: {filename}]\n{description}")
                processed_names.append(filename)
            else:
                failed_names.append(filename)
        else:
            text = extract_text(file_bytes, ext)
            if not text:
                failed_names.append(filename)
                continue
            text = safe_trim(clean_text(text), MAX_CHARS_PER_FILE)
            file_sections.append(f"[File: {filename}]\n{text}")
            processed_names.append(filename)

    if not file_sections:
        return jsonify({
            "reply": (
                "I opened the file(s) but couldn't extract readable content. "
                "If any are scanned documents, try uploading them as PNG or JPG images instead."
            )
        })

    combined_context = safe_trim("\n\n".join(file_sections), MAX_TOTAL_FILE_CHARS)

    instruction = (
        user_question if user_question
        else "Review all the uploaded files together and summarise the key information clearly and concisely."
    )

    msgs = [
        {
            "role": "system",
            "content": (
                f"You are Pranox AI. The user has uploaded {len(processed_names)} file(s). "
                "Use the extracted content from ALL files below to answer the user. "
                "Cross-reference information across files where relevant and mention "
                "specific filenames when it helps clarity. Be clear, structured, and helpful."
            ),
        },
        {"role": "user", "content": f"Uploaded files:\n\n{combined_context}\n\nUser request: {instruction}"},
    ]

    reply = run_ai_with_effort(msgs, effort)
    if not reply:
        reply = "I read the files but had trouble generating a response. Please try again."
    reply = re.sub(r"\n{3,}", "\n\n", reply).strip()

    if failed_names:
        reply += f"\n\n_Note: couldn't read {', '.join(failed_names)}._"

    db = get_db()
    try:
        file_label = ", ".join(processed_names) if processed_names else ", ".join(failed_names)
        db.execute(
            "INSERT INTO chats(user_email,role,message) VALUES (?,?,?)",
            (user_email, "user", f"[Files: {file_label}] {user_question}")
        )
        db.execute(
            "INSERT INTO chats(user_email,role,message) VALUES (?,?,?)",
            (user_email, "assistant", reply)
        )
        db.commit()
    finally:
        db.close()

    return jsonify({"reply": reply})


# ═══════════════════════════════════════════════════════
#  RESUME DOWNLOAD
# ═══════════════════════════════════════════════════════

@app.route("/download_resume", methods=["POST"])
def download_resume():
    buffer = BytesIO()
    p      = canvas.Canvas(buffer, pagesize=A4)
    _, height = A4
    y = height - 50
    p.setFont("Helvetica-Bold", 14)
    p.drawString(50, y, "Resume")
    y -= 30
    p.setFont("Helvetica", 11)
    for line in request.form["resume"].split("\n"):
        if y < 60:
            p.showPage()
            y = height - 50
            p.setFont("Helvetica", 11)
        line = line.strip()
        if line:
            p.drawString(50, y, line[:100])
            y -= 18
        else:
            y -= 8
    p.save()
    buffer.seek(0)
    return send_file(
        buffer,
        as_attachment=True,
        download_name="pranox_resume.pdf",
        mimetype="application/pdf",
    )


# ═══════════════════════════════════════════════════════
#  IMAGE GENERATION — Cloudflare
# ═══════════════════════════════════════════════════════

@app.route("/api/image", methods=["POST"])
def generate_image():
    data   = request.get_json(force=True)
    prompt = data.get("prompt", "").strip()
    if not prompt:
        return jsonify({"reply": "Please provide an image description."}), 400
    token   = CLOUDFLARE_API_TOKEN.strip()
    account = CLOUDFLARE_ACCOUNT_ID.strip()
    if not token or not account:
        return jsonify({"reply": "Image generation is not configured."}), 500
    enhanced = f"ultra realistic, 4k, highly detailed, sharp focus, professional photography, {prompt}"
    url      = f"https://api.cloudflare.com/client/v4/accounts/{account}/ai/run/@cf/stabilityai/stable-diffusion-xl-base-1.0"
    headers  = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
    try:
        response = requests.post(
            url, headers=headers,
            json={"prompt": enhanced, "num_steps": 20},
            timeout=60,
        )
        if response.status_code == 200:
            ct = response.headers.get("Content-Type", "")
            if "image" in ct or len(response.content) > 1000:
                return response.content, 200, {
                    "Content-Type": "image/png",
                    "Cache-Control": "no-cache",
                }
            return jsonify({"reply": "Image generation returned unexpected data. Please try again."}), 500
        elif response.status_code == 401:
            return jsonify({"reply": "Image generation auth failed. Check your Cloudflare API token."}), 500
        elif response.status_code == 403:
            return jsonify({"reply": "Image generation permission denied. Verify your Cloudflare account and token."}), 500
        else:
            return jsonify({"reply": f"Image generation failed (error {response.status_code})."}), 500
    except requests.Timeout:
        return jsonify({"reply": "Image generation timed out. Please try again."}), 504
    except Exception as e:
        print("IMAGE ERROR:", e)
        return jsonify({"reply": "Server error during image generation."}), 500


# ═══════════════════════════════════════════════════════
#  DEBUG — test search without opening the chat UI
#  Usage: visit /api/debug-search?q=who+is+cm+of+bengal
# ═══════════════════════════════════════════════════════

@app.route("/api/debug-search")
def debug_search():
    query         = request.args.get("q", "who is the cm of west bengal")
    do_search, sq = ai_search_decision(query)
    results       = search_internet(sq) if do_search and sq else "(search not triggered)"
    return jsonify({
        "original_query":    query,
        "search_triggered":  do_search,
        "search_query_used": sq,
        "serper_key_set":    bool(os.getenv("SERPER_API_KEY")),
        "results_length":    len(results),
        "results_preview":   results[:600] if results else "NO RESULTS",
    })



# ═══════════════════════════════════════════════════════
#  DEEPSEARCH ROUTES
# ═══════════════════════════════════════════════════════

@app.route("/deepsearch")
def deepsearch_page():
    return render_template("deepsearch.html", user=session.get("user"))


@app.route("/api/deepsearch", methods=["POST"])
def api_deepsearch():
    data = request.get_json(force=True)
    question = data.get("question", "").strip()
    if not question:
        return jsonify({"error": "Please provide a research question."}), 400

    def generate():
        for chunk in run_deepsearch(question):
            yield chunk

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        }
    )


# ═══════════════════════════════════════════════════════
#  STATIC PAGES
# ═══════════════════════════════════════════════════════

@app.route("/privacy")
def privacy():
    return render_template("privacy.html")

@app.route("/terms")
def terms():
    return render_template("terms.html")

@app.route("/cookies")
def cookies():
    return render_template("cookies.html")

@app.route("/changelog")
def changelog():
    return render_template("changelog.html")

@app.route("/docs")
def docs():
    return render_template("docs.html")

@app.route("/blog")
def blog():
    return render_template("blog.html")

@app.route("/help")
def help_center():
    return render_template("help.html")

@app.route("/status")
def status_page():
    return render_template("status.html")

@app.route("/sitemap.xml")
def sitemap():
    return send_from_directory("static", "sitemap.xml")

@app.route("/robots.txt")
def robots():
    return send_from_directory("static", "robots.txt")


if __name__ == "__main__":
    port = int(os.getenv("PORT", 8000))
    app.run(debug=False, host="0.0.0.0", port=port)